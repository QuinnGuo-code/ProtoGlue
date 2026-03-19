"""PowerPoint report generator for ProtoGlue results.

Requires the optional ``python-pptx`` dependency:
    pip install "protoglue[pptx]"
"""

import datetime
from pathlib import Path
from glob import glob
import numpy as np

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False


def generate_ppt(
    fig_dir: Path,
    output_path: Path,
    n_clusters: int,
    labels: np.ndarray,
    sil_val: float,
    ari_val: float,
    nmi_val: float,
    edge_purity: float,
    DATA_DIR: Path = None,
    LATENT_DIM: int = 64,
    NUM_HEADS: int = 8,
    NUM_LAYERS: int = 2,
    SPATIAL_SCALES: list = None,
    LR: float = 5e-4,
    PRETRAIN_EPOCHS: int = 300,
    FINETUNE_EPOCHS: int = 600,
):
    """Generate a PowerPoint summary of ProtoGlue results.

    Parameters
    ----------
    fig_dir : Path
        Directory containing saved figure PNGs.
    output_path : Path
        Output .pptx file path.
    n_clusters : int
        Number of clusters.
    labels : ndarray
        Cluster labels.
    sil_val, ari_val, nmi_val, edge_purity : float
        Evaluation metrics.
    DATA_DIR : Path, optional
        Data directory (for display in title).
    LATENT_DIM, NUM_HEADS, NUM_LAYERS, LR, etc.
        Hyperparameters to display.
    """
    if not HAS_PPTX:
        print("python-pptx not installed. Skipping PPT generation.")
        print("Install with: pip install python-pptx Pillow")
        return

    if SPATIAL_SCALES is None:
        SPATIAL_SCALES = [3, 8]

    fig_dir = Path(fig_dir)
    output_path = Path(output_path)

    # Color scheme (academic dark blue theme)
    C_NAVY = RGBColor(0x1E, 0x27, 0x61)
    C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    C_LIGHT = RGBColor(0xE8, 0xEB, 0xF5)
    C_ACCENT = RGBColor(0x3C, 0x54, 0x88)

    # Find available figures
    def find_fig(pattern):
        matches = sorted(glob(str(fig_dir / f"*{pattern}*")))
        hires = [m for m in matches if "_hires" in m]
        if hires:
            return hires[0]
        png = [m for m in matches if m.endswith(".png") and "_hires" not in m]
        return png[0] if png else None

    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    def add_title_bar(slide, title_text):
        """Add a colored title bar to the top of a slide."""
        from pptx.enum.shapes import MSO_SHAPE
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
            prs.slide_width, Inches(0.9),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = C_NAVY
        shape.line.fill.background()
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.color.rgb = C_WHITE
        p.font.size = Pt(24)
        p.font.bold = True
        p.alignment = PP_ALIGN.LEFT
        tf.margin_left = Inches(0.4)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # --- Slide 1: Title ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_title_bar(slide, "ProtoGlue Analysis Report")
    dataset_name = DATA_DIR.name if DATA_DIR else "Dataset"
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(12), Inches(4))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = f"Dataset: {dataset_name}"
    p.font.size = Pt(20)
    p.font.color.rgb = C_NAVY
    p = tf.add_paragraph()
    p.text = f"Date: {datetime.datetime.now().strftime('%Y-%m-%d')}"
    p.font.size = Pt(14)
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = f"Clusters: {n_clusters}  |  Spots: {len(labels)}"
    p.font.size = Pt(16)
    p = tf.add_paragraph()
    metrics_text = f"Silhouette: {sil_val:.4f}  |  Edge Purity: {edge_purity:.4f}"
    if np.isfinite(ari_val):
        metrics_text += f"  |  ARI: {ari_val:.4f}  |  NMI: {nmi_val:.4f}"
    p.text = metrics_text
    p.font.size = Pt(14)
    p.font.color.rgb = C_ACCENT
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = (
        f"Architecture: latent_dim={LATENT_DIM}, heads={NUM_HEADS}, "
        f"layers={NUM_LAYERS}, scales={SPATIAL_SCALES}"
    )
    p.font.size = Pt(12)
    p = tf.add_paragraph()
    p.text = f"Training: lr={LR}, pretrain={PRETRAIN_EPOCHS}ep, finetune={FINETUNE_EPOCHS}ep"
    p.font.size = Pt(12)

    # --- Slide 2+: Figures ---
    slide_specs = [
        ("overview", "Spatial Domain Overview"),
        ("scale_attn", "Multi-Scale Attention Weights"),
        ("kscan", "K Selection Scan"),
        ("training", "Training Curves"),
        ("heatmap", "Marker Gene Heatmap"),
        ("dotplot", "Marker Gene Dot Plot"),
        ("density", "UMAP Density"),
        ("spatial_multiples", "Per-Cluster Spatial Distribution"),
        ("summary", "Summary Panel"),
    ]

    for pattern, title in slide_specs:
        fig_path = find_fig(pattern)
        if fig_path is None:
            continue
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_title_bar(slide, title)
        try:
            slide.shapes.add_picture(
                fig_path, Inches(0.3), Inches(1.1),
                width=Inches(12.7), height=Inches(6.0),
            )
        except Exception as e:
            print(f"Warning: could not add {fig_path}: {e}")

    # Save
    prs.save(str(output_path))
    print(f"PPT saved to {output_path}")
